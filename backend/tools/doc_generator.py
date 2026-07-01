import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGB


P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'


THEME_PALETTES = [
    {
        "name": "Midnight Indigo",
        "bg": (0x0B, 0x0E, 0x17),
        "card": (0x14, 0x1A, 0x2E),
        "title": (0x81, 0x7A, 0xFF),
        "text": (0xE2, 0xE8, 0xF0),
        "accent1": (0x6C, 0x63, 0xFF),
        "accent2": (0x00, 0xD2, 0xFF),
        "gradient_start": (0x1A, 0x1A, 0x3E),
        "gradient_end": (0x0B, 0x0E, 0x17),
    },
    {
        "name": "Crimson Night",
        "bg": (0x0F, 0x0A, 0x15),
        "card": (0x1C, 0x14, 0x24),
        "title": (0xFF, 0x6B, 0x6B),
        "text": (0xF0, 0xEA, 0xE0),
        "accent1": (0xFF, 0x4D, 0x6D),
        "accent2": (0xFF, 0xD9, 0x3D),
        "gradient_start": (0x2A, 0x0E, 0x1E),
        "gradient_end": (0x0F, 0x0A, 0x15),
    },
    {
        "name": "Ocean Teal",
        "bg": (0x0A, 0x14, 0x1E),
        "card": (0x12, 0x22, 0x33),
        "title": (0x4E, 0xC5, 0xC1),
        "text": (0xDA, 0xE8, 0xEA),
        "accent1": (0x38, 0xB2, 0xAC),
        "accent2": (0xE8, 0x6D, 0x82),
        "gradient_start": (0x0E, 0x28, 0x38),
        "gradient_end": (0x0A, 0x14, 0x1E),
    },
    {
        "name": "Royal Purple",
        "bg": (0x10, 0x08, 0x1E),
        "card": (0x1E, 0x12, 0x38),
        "title": (0xC0, 0x79, 0xEA),
        "text": (0xE8, 0xE0, 0xF0),
        "accent1": (0xA8, 0x5C, 0xF0),
        "accent2": (0xF0, 0x71, 0xA9),
        "gradient_start": (0x22, 0x0E, 0x3A),
        "gradient_end": (0x10, 0x08, 0x1E),
    },
    {
        "name": "Ember Gold",
        "bg": (0x12, 0x10, 0x0A),
        "card": (0x22, 0x1E, 0x14),
        "title": (0xED, 0xB6, 0x46),
        "text": (0xF0, 0xE8, 0xDA),
        "accent1": (0xF5, 0x9E, 0x0B),
        "accent2": (0xED, 0x64, 0x36),
        "gradient_start": (0x2A, 0x22, 0x10),
        "gradient_end": (0x12, 0x10, 0x0A),
    },
]

TRANSITIONS = ['fade', 'push', 'wipe', 'cover', 'split', 'wheel']


def _rgb(t):
    return RGBColor(t[0], t[1], t[2])


def _add_transition(slide, transition_type='fade', speed='med'):
    trans = etree.SubElement(slide._element, f'{{{P_NS}}}transition')
    trans.set('spd', speed)
    trans.set('advClick', '1')

    type_map = {
        'fade': 'fade',
        'push': 'push',
        'wipe': 'wipe',
        'cover': 'cover',
        'split': 'split',
        'wheel': 'wheel',
    }
    tag = type_map.get(transition_type, 'fade')
    etree.SubElement(trans, f'{{{P_NS}}}{tag}')


def _add_entrance_animation(slide, shape_ids, delay_between=300):
    timing = etree.SubElement(slide._element, f'{{{P_NS}}}timing')
    tnLst = etree.SubElement(timing, f'{{{P_NS}}}tnLst')
    par_root = etree.SubElement(tnLst, f'{{{P_NS}}}par')
    cTn_root = etree.SubElement(par_root, f'{{{P_NS}}}cTn')
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, f'{{{P_NS}}}childTnLst')
    seq = etree.SubElement(childTnLst_root, f'{{{P_NS}}}seq')
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, f'{{{P_NS}}}cTn')
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, f'{{{P_NS}}}childTnLst')

    anim_id = 3
    for i, sp_id in enumerate(shape_ids):
        par_click = etree.SubElement(childTnLst_seq, f'{{{P_NS}}}par')
        cTn_click = etree.SubElement(par_click, f'{{{P_NS}}}cTn')
        cTn_click.set('id', str(anim_id))
        cTn_click.set('fill', 'hold')
        anim_id += 1

        stCondLst = etree.SubElement(cTn_click, f'{{{P_NS}}}stCondLst')
        cond = etree.SubElement(stCondLst, f'{{{P_NS}}}cond')
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, f'{{{P_NS}}}childTnLst')

        par_anim = etree.SubElement(childTnLst_click, f'{{{P_NS}}}par')
        cTn_anim = etree.SubElement(par_anim, f'{{{P_NS}}}cTn')
        cTn_anim.set('id', str(anim_id))
        cTn_anim.set('fill', 'hold')
        anim_id += 1

        stCond2 = etree.SubElement(cTn_anim, f'{{{P_NS}}}stCondLst')
        cond2 = etree.SubElement(stCond2, f'{{{P_NS}}}cond')
        cond2.set('delay', str(i * delay_between))

        childTnLst_anim = etree.SubElement(cTn_anim, f'{{{P_NS}}}childTnLst')

        animEffect = etree.SubElement(childTnLst_anim, f'{{{P_NS}}}animEffect')
        animEffect.set('transition', 'in')
        animEffect.set('filter', 'fade')

        cBhvr = etree.SubElement(animEffect, f'{{{P_NS}}}cBhvr')
        cTn_effect = etree.SubElement(cBhvr, f'{{{P_NS}}}cTn')
        cTn_effect.set('id', str(anim_id))
        cTn_effect.set('dur', '500')
        anim_id += 1

        tgtEl = etree.SubElement(cBhvr, f'{{{P_NS}}}tgtEl')
        spTgt = etree.SubElement(tgtEl, f'{{{P_NS}}}spTgt')
        spTgt.set('spid', str(sp_id))

    prevCondLst = etree.SubElement(seq, f'{{{P_NS}}}prevCondLst')
    prevCond = etree.SubElement(prevCondLst, f'{{{P_NS}}}cond')
    prevCond.set('evt', 'onPrev')
    prevCond.set('delay', '0')
    prevTgt = etree.SubElement(prevCond, f'{{{P_NS}}}tgtEl')
    etree.SubElement(prevTgt, f'{{{P_NS}}}sldTgt')

    nextCondLst = etree.SubElement(seq, f'{{{P_NS}}}nextCondLst')
    nextCond = etree.SubElement(nextCondLst, f'{{{P_NS}}}cond')
    nextCond.set('evt', 'onNext')
    nextCond.set('delay', '0')
    nextTgt = etree.SubElement(nextCond, f'{{{P_NS}}}tgtEl')
    etree.SubElement(nextTgt, f'{{{P_NS}}}sldTgt')


def _set_slide_bg(slide, color_tuple):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_tuple)


def _add_gradient_bar(slide, color_tuple, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_tuple)
    shape.line.fill.background()
    return shape


def _add_glow_circle(slide, color_tuple, left, top, size, alpha=30):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_tuple)
    shape.line.fill.background()

    sp_elem = shape._element
    srgb = sp_elem.find(f'.//{{{A_NS}}}srgbClr')
    if srgb is not None:
        alpha_elem = etree.SubElement(srgb, f'{{{A_NS}}}alpha')
        alpha_elem.set('val', str(alpha * 1000))
    return shape


def generate_pptx(title, slides_data, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = random.choice(THEME_PALETTES)
    transition_type = random.choice(TRANSITIONS)

    # ── TITLE SLIDE ──
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(title_slide, theme["bg"])

    _add_glow_circle(title_slide, theme["accent1"], Inches(-1), Inches(-1), Inches(5), alpha=8)
    _add_glow_circle(title_slide, theme["accent2"], Inches(9), Inches(4), Inches(6), alpha=6)

    _add_gradient_bar(title_slide, theme["accent1"], Inches(0), Inches(3.4), Inches(13.333), Inches(0.06))

    title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme["title"])
    p.alignment = PP_ALIGN.LEFT
    p.font.letter_spacing = Pt(2)

    sub_p = tf.add_paragraph()
    sub_p.text = "Generated by DevOps Concierge Agent"
    sub_p.font.size = Pt(16)
    sub_p.font.color.rgb = _rgb(theme["accent2"])
    sub_p.alignment = PP_ALIGN.LEFT
    sub_p.space_before = Pt(16)

    _add_gradient_bar(title_slide, theme["accent1"], Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

    shape_ids = [title_box.shape_id]
    _add_entrance_animation(title_slide, shape_ids)
    _add_transition(title_slide, transition_type)

    # ── CONTENT SLIDES ──
    for i, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, theme["bg"])

        _add_glow_circle(slide, theme["accent1"], Inches(10), Inches(-2), Inches(5), alpha=5)

        accent_bar = _add_gradient_bar(
            slide, theme["accent1"],
            Inches(0.8), Inches(0.6),
            Inches(0.07), Inches(0.8)
        )

        slide_title = slide_info.get("title", "")
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(10), Inches(1))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = slide_title
        tp.font.size = Pt(34)
        tp.font.bold = True
        tp.font.color.rgb = _rgb(theme["title"])

        _add_gradient_bar(slide, theme["accent2"], Inches(1.2), Inches(1.5), Inches(2), Inches(0.04))

        slide_num_box = slide.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
        snf = slide_num_box.text_frame
        snp = snf.paragraphs[0]
        snp.text = f"{i + 1:02d}"
        snp.font.size = Pt(14)
        snp.font.color.rgb = _rgb(theme["accent2"])
        snp.alignment = PP_ALIGN.RIGHT

        content_text = slide_info.get("content", "")
        bullets = slide_info.get("bullets", [])

        body_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(4.8))
        btf = body_box.text_frame
        btf.word_wrap = True

        if content_text:
            cp = btf.paragraphs[0]
            cp.text = content_text
            cp.font.size = Pt(18)
            cp.font.color.rgb = _rgb(theme["text"])
            cp.space_after = Pt(16)
            cp.line_spacing = Pt(28)

        for j, bullet in enumerate(bullets):
            bp = btf.add_paragraph() if (content_text or j > 0) else btf.paragraphs[0]
            bp.text = f"▸  {bullet}"
            bp.font.size = Pt(17)
            bp.font.color.rgb = _rgb(theme["text"])
            bp.space_before = Pt(10)
            bp.space_after = Pt(6)
            bp.line_spacing = Pt(26)

        _add_gradient_bar(slide, theme["accent1"], Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

        anim_shapes = [title_box.shape_id, body_box.shape_id]
        _add_entrance_animation(slide, anim_shapes, delay_between=400)
        _add_transition(slide, transition_type)

    output_path = os.path.join(output_dir, f"{title.replace(' ', '_')}.pptx")
    prs.save(output_path)
    return {"success": True, "file": output_path, "theme": theme["name"], "slides_count": len(slides_data) + 1}


def generate_docx(title, sections, output_dir):
    import os
    from datetime import datetime
    from docx import Document
    from docx.shared import Pt as DocPt, RGBColor as DocRGB, Inches as DocInches
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    os.makedirs(output_dir, exist_ok=True)
    doc = Document()

    # Configure 1-inch margins
    for section in doc.sections:
        section.top_margin = DocInches(1)
        section.bottom_margin = DocInches(1)
        section.left_margin = DocInches(1)
        section.right_margin = DocInches(1)

    # Helper to set cell background
    def set_cell_background(cell, fill_hex):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
        tcPr.append(shd)

    # Helper to set callout thick left border
    def set_callout_borders(cell, border_color_hex="4F46E5"):
        tcPr = cell._tc.get_or_add_tcPr()
        tcBorders = parse_xml(f'''
            <w:tcBorders {nsdecls("w")}>
                <w:top w:val="none"/>
                <w:left w:val="single" w:sz="36" w:space="0" w:color="{border_color_hex}"/>
                <w:bottom w:val="none"/>
                <w:right w:val="none"/>
            </w:tcBorders>
        ''')
        tcPr.append(tcBorders)

    # Configure base styles
    style_normal = doc.styles['Normal']
    style_normal.font.name = 'Segoe UI'
    style_normal.font.size = DocPt(11)
    style_normal.font.color.rgb = DocRGB(0x33, 0x41, 0x55) # Slate gray body
    style_normal.paragraph_format.line_spacing = 1.15
    style_normal.paragraph_format.space_after = DocPt(6)

    # ── 1. GORGEOUS COVER PAGE ──
    # Large space at the top
    top_spacer = doc.add_paragraph()
    top_spacer.paragraph_format.space_before = DocPt(100)

    # Document Category / Accent Tag
    tag_p = doc.add_paragraph()
    tag_run = tag_p.add_run("TECHNICAL SPECIFICATION & HANDOVER")
    tag_run.font.name = 'Segoe UI'
    tag_run.font.size = DocPt(10)
    tag_run.font.bold = True
    tag_run.font.color.rgb = DocRGB(0x4F, 0x46, 0xE5) # Accent Blue
    tag_p.paragraph_format.space_after = DocPt(12)

    # Main Title
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(title.upper())
    title_run.font.name = 'Segoe UI'
    title_run.font.size = DocPt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = DocRGB(0x0F, 0x17, 0x2A) # Slate 900 (Almost Black)
    title_p.paragraph_format.space_after = DocPt(6)

    # Subtitle / Description
    subtitle_p = doc.add_paragraph()
    sub_run = subtitle_p.add_run("Automated deployment architecture, database configuration, and system handover specifications generated by your DevOps Concierge.")
    sub_run.font.name = 'Segoe UI'
    sub_run.font.size = DocPt(12)
    sub_run.font.color.rgb = DocRGB(0x64, 0x74, 0x8B) # Slate 500
    subtitle_p.paragraph_format.space_after = DocPt(40)

    # Horizontal Accent Line
    line_table = doc.add_table(rows=1, cols=1)
    line_cell = line_table.rows[0].cells[0]
    set_cell_background(line_cell, "4F46E5") # Blue accent line
    line_table.rows[0].height = DocPt(4)

    # Spacer
    mid_spacer = doc.add_paragraph()
    mid_spacer.paragraph_format.space_before = DocPt(150)

    # Metadata Block
    meta_p = doc.add_paragraph()
    meta_run = meta_p.add_run(
        f"Prepared By: DevOps Concierge Agent v1.0\n"
        f"Date: {datetime.now().strftime('%B %d, %Y')}\n"
        f"Status: Active / Handover Ready"
    )
    meta_run.font.name = 'Segoe UI'
    meta_run.font.size = DocPt(10)
    meta_run.font.color.rgb = DocRGB(0x64, 0x74, 0x8B)
    meta_p.paragraph_format.line_spacing = 1.3

    # Page Break after Cover Page
    doc.add_page_break()

    # ── 2. SECTIONS RENDERING ──
    for sec_idx, section in enumerate(sections):
        # Heading 1
        h_text = section.get("heading", "")
        if h_text:
            hp = doc.add_paragraph()
            hrun = hp.add_run(h_text)
            hrun.font.name = 'Segoe UI'
            hrun.font.size = DocPt(18)
            hrun.font.bold = True
            hrun.font.color.rgb = DocRGB(0x0F, 0x17, 0x2A)
            hp.paragraph_format.space_before = DocPt(24)
            hp.paragraph_format.space_after = DocPt(8)
            hp.paragraph_format.keep_with_next = True

            # Underline separator for Heading 1
            line_table = doc.add_table(rows=1, cols=1)
            line_cell = line_table.rows[0].cells[0]
            set_cell_background(line_cell, "E2E8F0") # Subtle gray underline
            line_table.rows[0].height = DocPt(1)
            
            # Spacer after line
            space_p = doc.add_paragraph()
            space_p.paragraph_format.space_after = DocPt(8)

        # Paragraphs
        for para in section.get("paragraphs", []):
            para_lower = para.lower()
            if para_lower.startswith("note:") or para_lower.startswith("warning:") or para_lower.startswith("important:") or para_lower.startswith("tip:"):
                # Render as a gorgeous Callout Box!
                callout_table = doc.add_table(rows=1, cols=1)
                callout_cell = callout_table.rows[0].cells[0]
                
                bg_color = "F0FDF4" # Green bg for default/tip
                border_color = "10B981" # Green border
                if "warning" in para_lower or "important" in para_lower:
                    bg_color = "FEF2F2" # Red bg
                    border_color = "EF4444" # Red border
                elif "note" in para_lower:
                    bg_color = "F0F9FF" # Blue bg
                    border_color = "06B6D4" # Cyan/Blue border
                    
                set_cell_background(callout_cell, bg_color)
                set_callout_borders(callout_cell, border_color)
                
                cp = callout_cell.paragraphs[0]
                c_run = cp.add_run(para)
                c_run.font.name = 'Segoe UI'
                c_run.font.size = DocPt(10.5)
                c_run.font.italic = True
                c_run.font.color.rgb = DocRGB(0x1E, 0x29, 0x3B)
                cp.paragraph_format.space_after = DocPt(0)
                
                # Space after table
                post_table_p = doc.add_paragraph()
                post_table_p.paragraph_format.space_after = DocPt(8)
            else:
                p = doc.add_paragraph(para)

        # Bullet points
        bullets = section.get("bullets", [])
        if bullets:
            for item in bullets:
                p = doc.add_paragraph(style="List Bullet")
                p.paragraph_format.space_after = DocPt(3)
                run = p.add_run(item)
                run.font.name = 'Segoe UI'
                run.font.size = DocPt(11)
                run.font.color.rgb = DocRGB(0x33, 0x41, 0x55)

            # Space after list
            post_list_p = doc.add_paragraph()
            post_list_p.paragraph_format.space_after = DocPt(8)

    # ── 3. HEADERS & FOOTERS (For pages 2+) ──
    for section in doc.sections:
        section.different_first_page_header_footer = True
        
        # Header (Pages 2+)
        header = section.header
        header_p = header.paragraphs[0]
        header_p.text = f"{title} | DevOps Concierge Handover Specification"
        header_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        header_p.paragraph_format.space_after = DocPt(12)
        for run in header_p.runs:
            run.font.name = 'Segoe UI'
            run.font.size = DocPt(8.5)
            run.font.color.rgb = DocRGB(0x94, 0xA3, 0xB8) # Slate 400
            
        # Footer (Pages 2+)
        footer = section.footer
        footer_p = footer.paragraphs[0]
        footer_p.text = "Confidential - For Internal Use Only"
        footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in footer_p.runs:
            run.font.name = 'Segoe UI'
            run.font.size = DocPt(8.5)
            run.font.color.rgb = DocRGB(0x94, 0xA3, 0xB8)

    output_path = os.path.join(output_dir, f"{title.replace(' ', '_')}.docx")
    doc.save(output_path)
    return {"success": True, "file": output_path}


def generate_mermaid(diagram_type, nodes, edges):
    if diagram_type == "mindmap":
        lines = ["mindmap"]
        lines.append(f"  root(({nodes[0]}))")
        for node in nodes[1:]:
            lines.append(f"    {node}")
        return "\n".join(lines)

    if diagram_type == "flowchart":
        lines = ["flowchart TD"]
        for edge in edges:
            src, dst, label = edge.get("from"), edge.get("to"), edge.get("label", "")
            if label:
                lines.append(f"    {src} -->|{label}| {dst}")
            else:
                lines.append(f"    {src} --> {dst}")
        return "\n".join(lines)

    if diagram_type == "sequence":
        lines = ["sequenceDiagram"]
        for edge in edges:
            lines.append(f"    {edge['from']}->>{ edge['to']}: {edge.get('label', '')}")
        return "\n".join(lines)

    return ""


def generate_all_docs(title, content_data, output_dir):
    results = {}

    slides = content_data.get("slides", [])
    if slides:
        results["pptx"] = generate_pptx(title, slides, output_dir)

    sections = content_data.get("sections", [])
    if sections:
        results["docx"] = generate_docx(title, sections, output_dir)

    diagram = content_data.get("diagram")
    if diagram:
        mermaid_code = generate_mermaid(
            diagram.get("type", "flowchart"),
            diagram.get("nodes", []),
            diagram.get("edges", [])
        )
        mermaid_path = os.path.join(output_dir, f"{title.replace(' ', '_')}_diagram.md")
        with open(mermaid_path, "w") as f:
            f.write(f"```mermaid\n{mermaid_code}\n```\n")
        results["mermaid"] = {"success": True, "file": mermaid_path, "code": mermaid_code}

    return results
