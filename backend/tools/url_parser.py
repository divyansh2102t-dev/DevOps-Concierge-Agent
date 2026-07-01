import os
import glob
import shutil
import tempfile
import httpx
from bs4 import BeautifulSoup

def get_local_browser_profiles():
    """
    Scans the local Windows machine for Chrome, Edge, Brave, and Chrome Beta user data directories and their profiles.
    Returns a list of dictionaries with profile paths and labels.
    """
    username = os.environ.get("USERNAME") or os.getlogin()
    paths = []
    
    # Chrome
    chrome_base = f"C:/Users/{username}/AppData/Local/Google/Chrome/User Data"
    if os.path.exists(chrome_base):
        profiles = ["Default"] + [os.path.basename(p) for p in glob.glob(os.path.join(chrome_base, "Profile *"))]
        for p in profiles:
            p_path = os.path.join(chrome_base, p)
            if os.path.exists(p_path):
                paths.append({
                    "browser": "chrome",
                    "profile_name": p,
                    "path": p_path,
                    "label": f"Google Chrome - {p}"
                })
                
    # Edge
    edge_base = f"C:/Users/{username}/AppData/Local/Microsoft/Edge/User Data"
    if os.path.exists(edge_base):
        profiles = ["Default"] + [os.path.basename(p) for p in glob.glob(os.path.join(edge_base, "Profile *"))]
        for p in profiles:
            p_path = os.path.join(edge_base, p)
            if os.path.exists(p_path):
                paths.append({
                    "browser": "edge",
                    "profile_name": p,
                    "path": p_path,
                    "label": f"Microsoft Edge - {p}"
                })

    # Brave
    brave_base = f"C:/Users/{username}/AppData/Local/BraveSoftware/Brave-Browser/User Data"
    if os.path.exists(brave_base):
        profiles = ["Default"] + [os.path.basename(p) for p in glob.glob(os.path.join(brave_base, "Profile *"))]
        for p in profiles:
            p_path = os.path.join(brave_base, p)
            if os.path.exists(p_path):
                paths.append({
                    "browser": "brave",
                    "profile_name": p,
                    "path": p_path,
                    "label": f"Brave Browser - {p}"
                })

    # Chrome Beta
    beta_base = f"C:/Users/{username}/AppData/Local/Google/Chrome Beta/User Data"
    if os.path.exists(beta_base):
        profiles = ["Default"] + [os.path.basename(p) for p in glob.glob(os.path.join(beta_base, "Profile *"))]
        for p in profiles:
            p_path = os.path.join(beta_base, p)
            if os.path.exists(p_path):
                paths.append({
                    "browser": "chrome-beta",
                    "profile_name": p,
                    "path": p_path,
                    "label": f"Google Chrome Beta - {p}"
                })
                
    return paths


def robust_copy_file(src, dst):
    """Copies a file, opening it in read-only mode to bypass Windows sharing locks if the browser is running."""
    try:
        if not os.path.exists(src):
            return False
        os.makedirs(os.path.dirname(dst), exist_ok=True)
        with open(src, "rb") as f_in:
            with open(dst, "wb") as f_out:
                shutil.copyfileobj(f_in, f_out)
        return True
    except Exception:
        try:
            shutil.copy2(src, dst)
            return True
        except Exception:
            return False

def robust_copy_dir(src_dir, dst_dir):
    """Recursively copies essential directories using the robust read-only file copier."""
    if not os.path.exists(src_dir):
        return
    os.makedirs(dst_dir, exist_ok=True)
    for root, dirs, files in os.walk(src_dir):
        for file in files:
            src_file = os.path.join(root, file)
            rel_path = os.path.relpath(src_file, src_dir)
            dst_file = os.path.join(dst_dir, rel_path)
            robust_copy_file(src_file, dst_file)

def clone_profile_safely(profile_info):
    """
    Clones only the essential credential files and directories (Cookies, Local Storage, etc.)
    and the parent Local State file to a temp directory. Uses read-only binary streams to
    bypass OS-level file locks when the browser is running.
    """
    import uuid
    # Handle string path (backward compatibility)
    if isinstance(profile_info, str):
        profile_path = profile_info
        profile_name = os.path.basename(profile_path)
    else:
        profile_path = profile_info["path"]
        profile_name = profile_info["profile_name"]
        
    user_data_dir = os.path.dirname(profile_path)
    
    # Create a unique temp user data directory
    temp_user_data = os.path.join(tempfile.gettempdir(), f"devops_agent_userdata_{uuid.uuid4().hex[:8]}")
    os.makedirs(temp_user_data, exist_ok=True)
    
    # Copy Local State if it exists (crucial for cookie decryption!)
    local_state_path = os.path.join(user_data_dir, "Local State")
    if os.path.exists(local_state_path):
        robust_copy_file(local_state_path, os.path.join(temp_user_data, "Local State"))
            
    # Create the temp profile directory structure
    temp_profile_dir = os.path.join(temp_user_data, profile_name)
    os.makedirs(temp_profile_dir, exist_ok=True)
    
    # Define essential credential files and directories
    essential_files = [
        "Cookies",
        os.path.join("Network", "Cookies"),
        "Preferences",
        "Secure Preferences",
        "Login Data",
        "Web Data"
    ]
    
    essential_dirs = [
        "Local Storage",
        "Session Storage",
        "Sync Data"
    ]
    
    # Copy essential components
    for f in essential_files:
        robust_copy_file(os.path.join(profile_path, f), os.path.join(temp_profile_dir, f))
        
    for d in essential_dirs:
        robust_copy_dir(os.path.join(profile_path, d), os.path.join(temp_profile_dir, d))
        
    return {
        "user_data_dir": temp_user_data,
        "profile_directory": profile_name
    }


async def parse_url(url_or_path):
    """
    Unified resource parser. Automatically detects if the input is:
    1. A local file path (supports DOCX, PPTX, TXT, MD, Images, PDFs)
    2. A URL (attempts authenticated Playwright scraping using Chrome/Edge profiles, with HTTP fallback)
    """
    # ── CASE 1: LOCAL FILE PARSING ──
    if os.path.exists(url_or_path):
        ext = os.path.splitext(url_or_path)[1].lower()
        
        # Text/Markdown files
        if ext in (".txt", ".md", ".json", ".yaml", ".yml", ".ini", ".conf", ".env", ".xml"):
            try:
                with open(url_or_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                return {
                    "success": True,
                    "type": "text_file",
                    "title": os.path.basename(url_or_path),
                    "content": content[:100000] # Cap at 100k characters for model safety
                }
            except Exception as e:
                return {"success": False, "error": f"Failed to read text file: {str(e)}"}
                
        # Word Documents (.docx)
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(url_or_path)
                text = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        text.append(p.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            text.append(" | ".join(row_text))
                return {
                    "success": True,
                    "type": "word_document",
                    "title": os.path.basename(url_or_path),
                    "content": "\n".join(text)
                }
            except Exception as e:
                return {"success": False, "error": f"Failed to read Word document: {str(e)}"}

        # PowerPoint Presentations (.pptx)
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(url_or_path)
                text = []
                for idx, slide in enumerate(prs.slides):
                    text.append(f"\n--- Slide {idx + 1} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text.append(shape.text.strip())
                return {
                    "success": True,
                    "type": "powerpoint_presentation",
                    "title": os.path.basename(url_or_path),
                    "content": "\n".join(text)
                }
            except Exception as e:
                return {"success": False, "error": f"Failed to read PowerPoint presentation: {str(e)}"}

        # Images (.png, .jpg, .jpeg, .webp, .gif)
        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".gif"):
            try:
                from PIL import Image
                with Image.open(url_or_path) as img:
                    width, height = img.size
                    img_format = img.format
                return {
                    "success": True,
                    "type": "image",
                    "title": os.path.basename(url_or_path),
                    "content": f"[Local Image File]\nFormat: {img_format}\nDimensions: {width}x{height} pixels\nPath: {url_or_path}\n(Note: This image can be analyzed using multimodal vision capabilities)"
                }
            except Exception as e:
                return {"success": False, "error": f"Failed to read image metadata: {str(e)}"}

        # PDFs (.pdf)
        elif ext == ".pdf":
            try:
                # Attempt to extract text if pypdf is installed
                import pypdf
                reader = pypdf.PdfReader(url_or_path)
                text = []
                for idx, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text.append(f"--- Page {idx + 1} ---\n{page_text}")
                return {
                    "success": True,
                    "type": "pdf_document",
                    "title": os.path.basename(url_or_path),
                    "content": "\n".join(text)
                }
            except ImportError:
                return {
                    "success": False,
                    "error": "PyPDF library is not installed. Please run: 'pip install pypdf' to enable PDF text parsing."
                }
            except Exception as e:
                return {"success": False, "error": f"Failed to parse PDF: {str(e)}"}

        else:
            return {"success": False, "error": f"Unsupported local file extension: {ext}"}

    # ── CASE 2: URL SCRAPING ──
    # Check if playwright is installed for authenticated/SPA browser scraping
    playwright_available = False
    try:
        from playwright.async_api import async_playwright
        playwright_available = True
    except ImportError:
        pass

    # If Playwright is installed, try to scrape using Chrome/Edge profiles to inherit logins!
    if playwright_available:
        profiles = get_local_browser_profiles()
        if profiles:
            # We will try the first available profile (usually Google Chrome Default is the main one)
            selected_profile = profiles[0]
            cloned_info = clone_profile_safely(selected_profile)
            
            from playwright.async_api import async_playwright
            async with async_playwright() as p:
                try:
                    context = await p.chromium.launch_persistent_context(
                        user_data_dir=cloned_info["user_data_dir"],
                        headless=True,
                        args=[
                            f"--profile-directory={cloned_info['profile_directory']}",
                            "--disable-gpu",
                            "--no-sandbox"
                        ]
                    )
                    page = await context.new_page()
                    
                    # Navigate and wait for SPA JS rendering (networkidle)
                    await page.goto(url_or_path, timeout=40000, wait_until="networkidle")
                    await page.wait_for_timeout(3000) # Buffer to load chat messages
                    
                    title = await page.title()
                    content = await page.evaluate("() => document.body.innerText")
                    await context.close()
                    
                    # Cleanup cloned temp profile
                    try:
                        shutil.rmtree(cloned_info["user_data_dir"])
                    except Exception:
                        pass
                        
                    return {
                        "success": True,
                        "type": "web_page_browser",
                        "title": title,
                        "content": content[:50000],
                        "url": url_or_path,
                        "browser_used": selected_profile["label"]
                    }
                except Exception as e:
                    # Fallback to standard HTTP if Playwright profile launch fails
                    print(f"Playwright profile scraping failed, falling back to HTTP: {e}")
                    try:
                        shutil.rmtree(cloned_info["user_data_dir"])
                    except Exception:
                        pass

    # ── FALLBACK CASE: STANDARD HTTP SCRAPING ──
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=20) as client:
            response = await client.get(url_or_path)
            response.raise_for_status()
    except Exception as e:
        return {
            "success": False, 
            "error": f"HTTP fetch failed: {str(e)}" + 
                     ("\n(Tip: To scrape private/SPA pages like ChatGPT, run: 'pip install playwright && playwright install chromium')" if not playwright_available else "")
        }

    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title else ""
    meta_desc = ""
    meta_tag = soup.find("meta", attrs={"name": "description"})
    if meta_tag:
        meta_desc = meta_tag.get("content", "")

    text = soup.get_text(separator="\n", strip=True)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    clean_text = "\n".join(lines[:300])

    tip_message = ""
    if not playwright_available:
        tip_message = "\n\n*(💡 Tip: To scrape authenticated chats (e.g. ChatGPT) or dynamic SPAs, run: 'pip install playwright && playwright install chromium' in your terminal. The agent will then automatically inherit your local browser login sessions!)*"

    return {
        "success": True,
        "type": "web_page_http",
        "title": title,
        "description": meta_desc,
        "content": clean_text + tip_message,
        "url": str(response.url)
    }
